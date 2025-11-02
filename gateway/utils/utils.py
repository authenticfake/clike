from __future__ import annotations
import json
import os,  logging
import httpx
import io
import base64
import mimetypes
from typing import List, Dict, Optional
import docx
from pdfminer.high_level import extract_text
import openpyxl  # .xlsx
import xlrd  # .xls (legacy)
from pyxlsb import open_workbook as open_xlsb  # .xlsb (optional)
from pptx import Presentation  # .pptx




INLINE_MAX_FILE_KB   = int(os.getenv("INLINE_MAX_FILE_KB", "64"))
INLINE_MAX_TOTAL_KB  = int(os.getenv("INLINE_MAX_TOTAL_KB", "256"))
RAG_SIZE_THRESHOLD_KB = int(os.getenv("RAG_SIZE_THRESHOLD_KB", "64"))
RAG_TOP_K            = int(os.getenv("RAG_TOP_K", "12"))
log = logging.getLogger("gateway.utils")

# ===== RAG hooks (best-effort; non bloccanti) =====
def _rag_project_id(body: dict) -> str:
    pid = (body or {}).get("project_id")
    if isinstance(pid, str) and pid.strip():
        return pid.strip()
    return "default"

RAG_TOP_K = int(os.getenv("RAG_TOP_K", "8"))
def _rag_base_url() -> str:
    # es.: "http://localhost:8080/v1/rag"
    base =  os.getenv("RAG_BASE_URL", "http://orchestrator:8080/v1/rag")
    return base.rstrip("/")



async def rag_index_items(project_id: str, items: list[dict]):
    # Optional server-side index; we prefer client-side, but keep for completeness.
    if not items:
        return
    payload = {"project_id": project_id, "items": []}
    for it in (items or []):
        p = (it.get("path") or "").strip()
        t = (it.get("text") or "").strip()
        if p and t:
            payload["items"].append({"path": p, "text": t})
    if not payload["items"]:
        return
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            await client.post(f"{_rag_base_url()}/index", json=payload)
    except Exception as e:
        log.warning("rag_index_items failed: %s", e)
        raise e


async def rag_fetch(project_id: str, paths: list[str]):
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(f"{_rag_base_url()}/fetch_by_paths",
                                  json={"project_id": project_id,
                                        "paths": paths,
                                        "max_chars_per_doc": 200000})
            r.raise_for_status()
            data = r.json() or {}
            return data
        
    except Exception as e:
        log.warning("rag_fetch failed: %s", e)
        return []


async def rag_query(project_id: str, query: str, top_k: int = None):
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(f"{_rag_base_url()}/search",
                                  json={"project_id": project_id,
                                        "query": query or "",
                                        "top_k": int(top_k or RAG_TOP_K)})
            r.raise_for_status()
            data = r.json() or {}
            return data.get("hits") or []
    except Exception as e:
        log.warning("rag_query failed: %s", e)
        return []
    
def _kb(n_bytes: int) -> int:
    try: return int(n_bytes) // 1024
    except: return 0
    
async def decide_inline_or_rag(attachments: list[dict]) -> tuple[list[dict], list[dict]]:
    """
    Allineato all'estensione (partitionAttachments):
      - inline se a.content o a.bytes_b64 presenti
      - altrimenti, se a.path presente => RAG by path
      - altrimenti ignora (log warning)
    Niente soglie di size, niente budget qui (per coerenza end-to-end).
    """
    inline, rag = [], []
    if not attachments:
        return inline, rag

    for a in attachments:
        if not isinstance(a, dict):
            continue

        name   = a.get("name") or a.get("path") or "file"
        path   = a.get("path")
        origin = a.get("origin") or a.get("source")  # normalizza
        content    = a.get("content")
        bytes_b64  = a.get("bytes_b64")

        # Nota: evitiamo di loggare la base64 (solo boolean), per non intasare i log
        log.info("decide inline or rag: %s",
                 json.dumps({
                     "name": name,
                     "has_content": bool(content),
                     "has_bytes_b64": bool(bytes_b64),
                     "path": path,
                     "origin": origin
                 }, ensure_ascii=False))

        if content or bytes_b64:
            # Inline esattamente come fa l’estensione
            if bytes_b64:
                raw = _b64_to_bytes(bytes_b64)
                if raw:
                    ext = _ext_from_path(path)
                    if ext == ".pdf":
                        log.info("inline: PDF")
                        txt = _extract_text_from_pdf_bytes(raw)
                    elif ext == ".docx":
                        log.info("inline DOCX")
                        txt = _extract_text_from_docx_bytes(raw)
                    elif ext == ".xlsx":
                        log.info("inline: XLSX")
                        txt = _extract_text_from_xlsx_bytes(raw)
                    elif ext == ".xls":
                        log.info("inline: XLS (legacy)")
                        txt = _extract_text_from_xls_bytes(raw)
                    elif ext == ".xlsb":
                        log.info("inline: XLSB")
                        txt = _extract_text_from_xlsb_bytes(raw)
                    elif ext == ".pptx":
                        log.info("inline: PPTX")
                        txt = _extract_text_from_pptx_bytes(raw)
                    else:
                        # fallback: se è testo “grezzo” o sconosciuto, prova a decodare come utf-8
                        try:
                            txt = raw.decode("utf-8", errors="ignore")
                        except Exception:
                            txt = ""
                  
                    log.info("inline file %s -> %d chars", path, len(txt))
                    content = txt

            inline.append({
                "name": name,
                "path": path,          # opzionale (può servire per tracciabilità)
                "content": content,    # può essere None
                "bytes_b64": bytes_b64,# può essere None
                "origin": origin
            })
        elif path:
            # RAG by path, minimale (non inoltriamo bytes_b64 per non gonfiare la payload)
            rag.append({
                "name": name,
                "path": path,
                "origin": origin
            })
        else:
            log.warning("Attachment senza content/bytes_b64 e senza path: ignorato: %s", name)

    log.info("attachments routing %s",
             json.dumps({"inline": len(inline), "rag": len(rag)}, ensure_ascii=False))
    return inline, rag

def _b64_to_bytes(b64: Optional[str]) -> Optional[bytes]:
    if not isinstance(b64, str) or not b64:
        return None
    try:
        # strip data URLs if present
        if b64.startswith("data:"):
            head, _, rest = b64.partition(",")
            b64 = rest
        return base64.b64decode(b64, validate=False)
    except Exception:
        return None

def _ext_from_path(p: str) -> str:
    try:
        return os.path.splitext((p or "").strip())[1].lower()
    except Exception:
        return ""
# --- ADD: extraction helpers ---
def _extract_text_from_pdf_bytes(raw: bytes) -> str:
    # pdfminer.six
    try:
        return extract_text(io.BytesIO(raw)) or ""
    except Exception as e:
        log.warning("PDF extract failed: %s", e)
        return ""
    


def _extract_text_from_docx_bytes(raw: bytes) -> str:
    # python-docx
    try:
        doc = docx.Document(io.BytesIO(raw))
        parts = []
        # paragraphs
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
        # tables (cells)
        for tbl in getattr(doc, "tables", []):
            for row in tbl.rows:
                for cell in row.cells:
                    t = (cell.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        log.warning("DOCX extract failed: %s", e)
        return ""
    
def _extract_text_from_xlsx_bytes(raw: bytes) -> str:
    if not openpyxl:
        log.warning("openpyxl non disponibile: skip xlsx")
        return ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLSX extract failed: %s", e)
        return ""

def _extract_text_from_xls_bytes(raw: bytes) -> str:
    # Richiede xlrd>=2.0 (legge solo .xls)
    if not xlrd:
        log.warning("xlrd non disponibile: skip xls")
        return ""
    try:
        book = xlrd.open_workbook(file_contents=raw)
        parts = []
        for si in range(book.nsheets):
            sh = book.sheet_by_index(si)
            parts.append(f"# Sheet: {sh.name}")
            for r in range(sh.nrows):
                row = [str(sh.cell_value(r, c)) for c in range(sh.ncols)]
                line = "\t".join(row).strip()
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLS extract failed: %s", e)
        return ""

def _extract_text_from_xlsb_bytes(raw: bytes) -> str:
    if not open_xlsb:
        log.warning("pyxlsb non disponibile: skip xlsb")
        return ""
    try:
        parts = []
        with open_xlsb(io.BytesIO(raw)) as wb:
            for sheet_name in wb.sheets:
                parts.append(f"# Sheet: {sheet_name}")
                with wb.get_sheet(sheet_name) as sh:
                    for row in sh.rows():
                        vals = [str(c.v) if c.v is not None else "" for c in row]
                        line = "\t".join(vals).strip()
                        if line:
                            parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLSB extract failed: %s", e)
        return ""

def _extract_text_from_pptx_bytes(raw: bytes) -> str:
    if not Presentation:
        log.warning("python-pptx non disponibile: skip pptx")
        return ""
    try:
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text or "" for run in para.runs).strip()
                        if text:
                            parts.append(text)
                elif hasattr(shape, "text") and shape.text:
                    t = (shape.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        log.warning("PPTX extract failed: %s", e)
        return ""


# --- RAG (http-based) collector ------------------------------------------------
async def fetch_rag_materials(project_id: str, queries: list[str] | None,  top_k: int | None = None) -> list[dict]:
    log.info("fetch_rag_materials: %s", json.dumps({"project_id": project_id, "queries": queries, "top_k": top_k}, ensure_ascii=False))

    data = await rag_fetch(project_id, queries)
    
    log.info("fetch_rag_materials: %s", json.dumps(data, ensure_ascii=False))
    return data

async def collect_rag_materials_http(
    project_id: str ,
    queries: list[str] | None,
    core_blobs: dict | None,
    top_k: int | None = None,
) -> list[dict]:
    """
    Interroga il servizio RAG via utils.rag_query() per ogni query.
    Ritorna: [{title, text, source}] senza duplicati, cap ~12 estratti.
    """
    pid = (project_id or "default").strip()
    qlist: list[str] = []

    # 1) Se non arrivano query, creale in base a path noti e heading dei core_blobs
    if not queries:
        # path-based (gli stessi che l'estensione indicizza)
        # qlist.extend([
        #     "path:docs/harper/README.md",
        #     "path:docs/harper/SPEC.md",
        #     "path:docs/harper/PLAN.md",
        #     "path:docs/harper/plan.json",
        #     "path:docs/harper/KIT.md",
        #     "path:docs/harper/",
        #     "path:src/",
        # ])
        qlist.extend([
            "path:src/",
        ])
        # heading-based (prima linea # ... di SPEC/PLAN se presenti nei core_blobs)
        # for key in ("SPEC.md", "PLAN.md"):
        #     txt = (core_blobs or {}).get(key, "") or ""
        #     for ln in txt.splitlines():
        #         if ln.strip().startswith("#"):
        #             qlist.append(ln.strip("# ").strip())
        #             break
    else:
        qlist = list(queries)

    # 2) esegui le query
    materials: list[dict] = []
    seen = set()
    cap = min(int(top_k or RAG_TOP_K), 120)
    for q in qlist[:8]:               # massimo 8 query
        hits = await rag_query(pid, q, top_k=cap)
        for h in (hits or []):
            path = (h.get("path") or "").strip()
            text = (h.get("text") or "").strip()
            if not text:
                continue
            key = (path, h.get("chunk"))
            if key in seen:
                continue
            seen.add(key)
            title = f"{path or 'doc'}#{h.get('chunk',0)}"
            if "score" in h:
                try:
                    title += f" (score={float(h['score']):.3f})"
                except Exception:
                    pass
            item = {"title": title, "text": text, "source": "rag"}
            log.info("RAG item: %s", item)
            materials.append(item)
            if len(materials) >= cap:
                break
        if len(materials) >= cap:
            break
    return materials

