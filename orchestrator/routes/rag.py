# orchestrator/routes/rag.py
from __future__ import annotations
from fastapi import APIRouter, HTTPException
from pydantic import BaseModel
from pydantic import BaseModel

from typing import List, Dict, Any, Optional
import logging
import io, os, base64
from pdfminer.high_level import extract_text  # import INSIDE to avoid import warning at module load

import docx
import pydantic
import openpyxl  # .xlsx
import xlrd  # .xls (legacy)
from pyxlsb import open_workbook as open_xlsb  # .xlsb (optional)
from pptx import Presentation  # .pptx
# Detect Pydantic major version once
try:
    _PYD_VER = int((getattr(pydantic, "__version__", "1.0.0").split(".")[0]) or "1")
except Exception:
    _PYD_VER = 1

# A small base that works on both Pydantic v1 and v2
if _PYD_VER >= 2:
    try:
        from pydantic import ConfigDict  # type: ignore
    except Exception:
        ConfigDict = dict  # fallback for type checkers
    class RagBase(BaseModel):
        model_config = ConfigDict(extra="ignore")  # v2 way
else:
    class RagBase(BaseModel):
        class Config:  # v1 way
            extra = "ignore"


from services.rag_store import RagStore

log = logging.getLogger("router.rag")

router = APIRouter()

class RagIndexItem(RagBase):
    path: str
    text: Optional[str] = None
    bytes_b64: Optional[str] = None

class RagIndexRequest(RagBase):
    project_id: str
    items: List[RagIndexItem]

class RagSearchRequest(RagBase):
    project_id: str
    query: str
    top_k: int = 8

class RagPurgeRequest(BaseModel):
    project_id: str
    path_prefix: Optional[str] = None

# --- NEW: fetch models ---
class RagFetchRequest(RagBase):
    project_id: str
    # Se indicati, limita il fetch a questi path (match "starts with" case-insensitive)
    paths: Optional[List[str]] = None
    # In alternativa/aggiunta, filtra per prefisso
    path_prefix: Optional[str] = None
    # Quanti documenti (path) restituire al massimo
    limit_docs: int = 20
    # Quanti caratteri massimi per documento aggregato (per prompt budget)
    max_chars_per_doc: int = 4000
    # Quanti chunk pescare dallo store (esageriamo: 5x documents)
    search_top_k: int = 100

class RagFetchByPathsRequest(RagBase):
    project_id: str
    paths: List[str]
    max_chars_per_doc: int = 4000
    search_top_k: int = 100
    
def _path_matches(p: str, paths: Optional[List[str]], prefix: Optional[str]) -> bool:
    p_norm = (p or "").strip()
    if not p_norm:
        return False
    p_low = p_norm.lower()
    if prefix and p_low.startswith(prefix.lower()):
        return True
    if paths:
        for want in paths:
            w = (want or "").strip()
            if not w:
                continue
            # match "starts with" per robustezza su path normalizzati
            if p_low.startswith(w.lower()):
                return True
    # se non sono imposti paths/prefix, accetta tutti
    return (paths is None and prefix is None)

def _aggregate_hits_by_path(
    hits: List[Dict[str, Any]],
    max_chars_per_doc: int,
    limit_docs: int,
    paths: Optional[List[str]],
    prefix: Optional[str],
) -> List[Dict[str, Any]]:
    """
    Raggruppa i risultati per 'path' e concatena i testi finché non supera max_chars_per_doc.
    Ritorna una lista di {path, text, chunks:int}.
    """
    buckets: Dict[str, Dict[str, Any]] = {}
    for h in (hits or []):
        p = (h.get("path") or "").strip()
        t = (h.get("text") or "").strip()
        if not p or not t:
            continue
        if not _path_matches(p, paths, prefix):
            continue
        b = buckets.get(p)
        if not b:
            b = {"path": p, "text": "", "chunks": 0}
            buckets[p] = b
        # Accumula rispettando il budget caratteri
        remaining = max_chars_per_doc - len(b["text"])
        if remaining <= 0:
            continue
        # +1 riga separatrice per chiarezza
        piece = (("\n" if b["text"] else "") + t)[:remaining]
        if piece:
            b["text"] += piece
            b["chunks"] += 1

    # Ordina per path (stabile) e limita la quantità di documenti
    ordered = list(buckets.values())
    ordered.sort(key=lambda x: x["path"])
    return ordered[: max(1, limit_docs)]


def _b64_to_bytes(b64: Optional[str]) -> Optional[bytes]:
    if not isinstance(b64, str) or not b64:
        return None
    try:
        # strip data URLs if present
        if b64.startswith("data:"):
            head, _, rest = b64.partition(",")
            b64 = rest
        return base64.b64decode(b64, validate=False)
    except Exception:
        return None

def _ext_from_path(p: str) -> str:
    try:
        return os.path.splitext((p or "").strip())[1].lower()
    except Exception:
        return ""
# --- ADD: extraction helpers ---
def _extract_text_from_pdf_bytes(raw: bytes) -> str:
    # pdfminer.six
    try:
        return extract_text(io.BytesIO(raw)) or ""
    except Exception as e:
        log.warning("PDF extract failed: %s", e)
        return ""
    


def _extract_text_from_docx_bytes(raw: bytes) -> str:
    # python-docx
    try:
        doc = docx.Document(io.BytesIO(raw))
        parts = []
        # paragraphs
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
        # tables (cells)
        for tbl in getattr(doc, "tables", []):
            for row in tbl.rows:
                for cell in row.cells:
                    t = (cell.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        log.warning("DOCX extract failed: %s", e)
        return ""
    
@router.post("/v1/rag/fetch")
async def rag_fetch(req: RagFetchRequest):
    """
    Ritorna il contenuto aggregato per documento (path) senza dover specificare una query utente.
    Usa internamente RagStore.search("", top_k=...) e filtra/accorpa lato router.
    """
    store = RagStore(project_id=req.project_id)

    # 1) Peschiamo tanti chunk neutrali (query vuota / "context")
    #    Nota: se il tuo RagStore non gestisce bene query vuota, prova con "context" o "*"
    query = ""
    raw_hits = await store.search(query, top_k=max(10, req.search_top_k))

    # 2) Aggrega per path rispettando i limiti
    docs = _aggregate_hits_by_path(
        hits=raw_hits,
        max_chars_per_doc=max(500, req.max_chars_per_doc),
        limit_docs=max(1, req.limit_docs),
        paths=req.paths,
        prefix=(req.path_prefix or None),
    )

    return {"docs": docs, "count": len(docs)}

def _extract_text_from_xlsx_bytes(raw: bytes) -> str:
    if not openpyxl:
        log.warning("openpyxl non disponibile: skip xlsx")
        return ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLSX extract failed: %s", e)
        return ""

def _extract_text_from_xls_bytes(raw: bytes) -> str:
    # Richiede xlrd>=2.0 (legge solo .xls)
    if not xlrd:
        log.warning("xlrd non disponibile: skip xls")
        return ""
    try:
        book = xlrd.open_workbook(file_contents=raw)
        parts = []
        for si in range(book.nsheets):
            sh = book.sheet_by_index(si)
            parts.append(f"# Sheet: {sh.name}")
            for r in range(sh.nrows):
                row = [str(sh.cell_value(r, c)) for c in range(sh.ncols)]
                line = "\t".join(row).strip()
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLS extract failed: %s", e)
        return ""

def _extract_text_from_xlsb_bytes(raw: bytes) -> str:
    if not open_xlsb:
        log.warning("pyxlsb non disponibile: skip xlsb")
        return ""
    try:
        parts = []
        with open_xlsb(io.BytesIO(raw)) as wb:
            for sheet_name in wb.sheets:
                parts.append(f"# Sheet: {sheet_name}")
                with wb.get_sheet(sheet_name) as sh:
                    for row in sh.rows():
                        vals = [str(c.v) if c.v is not None else "" for c in row]
                        line = "\t".join(vals).strip()
                        if line:
                            parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLSB extract failed: %s", e)
        return ""

def _extract_text_from_pptx_bytes(raw: bytes) -> str:
    if not Presentation:
        log.warning("python-pptx non disponibile: skip pptx")
        return ""
    try:
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text or "" for run in para.runs).strip()
                        if text:
                            parts.append(text)
                elif hasattr(shape, "text") and shape.text:
                    t = (shape.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        log.warning("PPTX extract failed: %s", e)
        return ""

@router.post("/v1/rag/fetch_by_paths")
async def rag_fetch_by_paths(req: RagFetchByPathsRequest):
    store = RagStore(project_id=req.project_id)
    log.info("RAG Store %s", store)
    # query neutra + filtro per path(s) lato router
    raw_hits = await store.search("", top_k=max(10, req.search_top_k))
    
    docs = _aggregate_hits_by_path(
        hits=raw_hits,
        max_chars_per_doc=max(20500, req.max_chars_per_doc),
        limit_docs=len(req.paths) if req.paths else 20,
        paths=req.paths,
        prefix=None,
    )
    return {"docs": docs, "count": len(docs)}

@router.post("/v1/rag/reindex")
async def rag_index(req: RagIndexRequest):
    return rag_index(req)

@router.post("/v1/rag/index")
async def rag_index(req: RagIndexRequest):
    # store = RagStore(project_id=req.project_id)
    # log.info("RAG Store %s", store)
    # log.info("RAG index %d items", len(req.items))
    # out = await store.index_texts([it.dict() for it in req.items])
    # log.info("RAG out %s", out)
    # if not out.get("ok"):
    #     raise HTTPException(500, detail=out.get("error","index failed"))
    # return out
   
    store = RagStore(project_id=req.project_id)
    log.info("RAG Store %s", store)
    log.info("RAG index %d items", len(req.items))

    # Costruisci docs normalizzati: sempre {"path":..., "text":...}
    docs = []
    for it in (req.items or []):
        p = (it.path or "").strip()
        txt = (it.text or "") if isinstance(it.text, str) else ""
        b64 = it.bytes_b64 or ""  # opzionale

        if not txt and b64:
            raw = _b64_to_bytes(b64)
            if raw:
                ext = _ext_from_path(p)
                if ext == ".pdf":
                    log.info("RAG index: PDF")
                    txt = _extract_text_from_pdf_bytes(raw)
                elif ext == ".docx":
                    log.info("RAG index: DOCX")
                    txt = _extract_text_from_docx_bytes(raw)
                elif ext == ".xlsx":
                    log.info("inline: XLSX")
                    txt = _extract_text_from_xlsx_bytes(raw)
                elif ext == ".xls":
                    log.info("inline: XLS (legacy)")
                    txt = _extract_text_from_xls_bytes(raw)
                elif ext == ".xlsb":
                    log.info("inline: XLSB")
                    txt = _extract_text_from_xlsb_bytes(raw)
                elif ext == ".pptx":
                    log.info("inline: PPTX")
                    txt = _extract_text_from_pptx_bytes(raw)
                else:
                    # fallback: se è testo “grezzo” o sconosciuto, prova a decodare come utf-8
                    try:
                        txt = raw.decode("utf-8", errors="ignore")
                    except Exception:
                        txt = ""
                
                log.info("RAG file %s -> %d chars", p, len(txt))

        if isinstance(txt, str) and txt.strip():
            docs.append({"path": p or "doc", "text": txt.strip()})

    if not docs:
        # nessun testo estraibile -> ok a vuoto (oppure alza 400 se preferisci)
        log.info("RAG index: no indexable docs")
        return {"ok": True, "count": 0}

    out = await store.index_texts(docs)
    log.info("RAG out %s", out)
    if not out.get("ok"):
        raise HTTPException(500, detail=out.get("error", "index failed"))
    return out


@router.post("/v1/rag/search")
async def rag_search(req: RagSearchRequest):
    store = RagStore(project_id=req.project_id)
    hits = await store.search(req.query, top_k=req.top_k)
    return {"hits": hits}

@router.post("/v1/rag/purge")
async def rag_purge(req: RagPurgeRequest):
    store = RagStore(project_id=req.project_id)
    out = await store.purge(req.path_prefix)
    if not out.get("ok"):
        raise HTTPException(500, detail=out.get("error","purge failed"))
    return out

# --- RAG merging: client chunks + server search (Qdrant) --------------------

def _chunk_map_from_client(rag_chunks: dict) -> dict:
    """Crea un dizionario {(name, idx) -> text} dai rag_chunks client."""
    cmap = {}
    for ch in (rag_chunks or []):
        name = (ch.get("name") or "").strip()
        idx  = ch.get("idx")
        txt  = (ch.get("text") or "").strip()
        if not name or idx is None or not txt:
            continue
        cmap[(name, int(idx))] = txt
    return cmap
