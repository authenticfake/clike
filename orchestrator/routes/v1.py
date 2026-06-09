# routes/v1.py
import os, json, logging, re, uuid, base64
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple
from fastapi import APIRouter, HTTPException, Request, Query
import httpx
from pydantic import BaseModel
from config import settings
from services import utils as su
from services.llm_client import call_gateway_chat, call_gateway_generate
from services.llm_contracts import resolve_llm_selection, load_catalog
import time as _time
from copy import deepcopy as _deepcopy
from services.rag_store import RagStore
import docx
from pdfminer.high_level import extract_text
import openpyxl  # .xlsx
import xlrd  # .xls (legacy)
from pyxlsb import open_workbook as open_xlsb  # .xlsb (optional)
from pptx import Presentation  # .pptx
# --- Generated root selection -------------------------------------------------
import uuid
from services.mode_contracts import normalize_mode_contract, validate_chat_contract, apply_generate_contract
from services.execution_policy import normalize_execution_preference

# splitter (alcune funzioni potrebbero non essere usate, ma manteniamo le import per compat)
from services.splitter import (
    infer_language,
    split_python_per_symbol,
    split_ts_per_symbol,
    apply_strategy,
)
def build_response_format_files_bundle() -> dict:
    """
    OpenAI structured output schema for a bundle of files.
    Strict schema: properties == required (no extras).
    Minimal: path, content, mime.
    """
    return {
        "type": "json_schema",
        "json_schema": {
            "name": "files_bundle_v1",
            "strict": True,
            "schema": {
                "type": "object",
                "properties": {
                    "files": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "path":    {"type": "string"},
                                "content": {"type": "string"},
                                "mime":    {"type": "string"},
                            },
                            "required": ["path", "content", "mime"],
                            "additionalProperties": False,
                        },
                        "minItems": 1
                    }
                },
                "required": ["files"],
                "additionalProperties": False,
            },
        },
    }

router = APIRouter(prefix="/v1")
log = logging.getLogger("orchestrator.v1")

INLINE_MAX_FILE_KB   = int(os.getenv("INLINE_MAX_FILE_KB", "64"))
INLINE_MAX_TOTAL_KB  = int(os.getenv("INLINE_MAX_TOTAL_KB", "256"))
RAG_SIZE_THRESHOLD_KB = int(os.getenv("RAG_SIZE_THRESHOLD_KB", "64"))
RAG_TOP_K            = int(os.getenv("RAG_TOP_K", "12"))



# --- Classification for src/doc buckets ---
CODE_EXTS = {
    ".py",".ts",".tsx",".js",".jsx",".go",".java",".c",".h",".cpp",".hpp",".cs",".rs",".kt",".swift",
    ".php",".rb",".pl",".r",".m",".scala",".sh",".ps1",".sql",".html",".css",".scss",".less",".xml",".xsl",
    ".json",".yaml",".yml",".toml",".ini",".gradle",".pom",".mdx",".vue",".svelte",".sol",".dart"
}
DOC_EXTS = {
    ".md",".rst",".txt",".adoc",".pdf",".doc",".docx",".ppt",".pptx",".odt",".rtf",".csv",".xlsx",".xls",".ipynb",
    ".mendixmodel",".mxmodel"
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp", ".bmp", ".tif", ".tiff"}
DATA_URL_RE = re.compile(r'data:(image/[\w\-\+\.]+);base64,([A-Za-z0-9+/=]+)')

def _get_cfg(name: str, default: str) -> str:
    """Legge prima da os.environ, poi da settings, altrimenti default."""
    v = os.getenv(name)
    if v is not None and str(v).strip():
        return str(v).strip()
    try:
        vv = getattr(settings, name, None)
        if vv is not None and str(vv).strip():
            return str(vv).strip()
    except Exception:
        pass
    return default

def _rag_base_url() -> str:
    # es.: "http://localhost:8080/v1/rag"
    base = _get_cfg("RAG_BASE_URL", "http://localhost:8080/v1/rag")
    return base.rstrip("/")

# --- RAG / attachments normalization helper ---------------------------------
def _normalize_context_from_body(body: dict) -> tuple[list[dict], list[dict], list[dict]]:
    """
    Return (inline_files, rag_files, attachments) normalized from request body.

    Accepted shapes:
      - inline_files / in_line_files: [{ "name"|"path", "content": "<text>" }]
      - rag_files: [{ "name"|"path", "path": "<abs-or-rel>", "bytes_b64": "<b64-optional>", "size": <int-optional> }]
      - attachments: VSCode-style attachment objects (will be auto-partitioned by _decide_inline_or_rag)

    We do NOT merge the legacy rag_paths/rag_inline here. That compatibility path
    is intentionally handled later and only if new-style inputs are empty.
    """
    if not isinstance(body, dict):
        return [], [], []

    inline_raw = body.get("in_line_files") or body.get("inline_files") or []
    rag_raw    = body.get("rag_files") or []
    atts_raw   = body.get("attachments") or []

    inline_files: list[dict] = []
    for item in inline_raw or []:
        if not isinstance(item, dict):
            continue
        name = (item.get("name") or item.get("path") or "file").strip()
        content = item.get("content")
        if isinstance(content, str) and content:
            inline_files.append({"name": name, "content": content})

    rag_files: list[dict] = []
    for item in rag_raw or []:
        if not isinstance(item, dict):
            continue
        name = (item.get("name") or item.get("path") or "").strip()
        path = (item.get("path") or "").strip()
        b64  = item.get("bytes_b64")
        size = item.get("size")
        rag_files.append({"name": name or (path or "file"), "path": path, "bytes_b64": b64, "size": size})

    attachments: list[dict] = []
    for item in atts_raw or []:
        if isinstance(item, dict):
            attachments.append(item)

    return inline_files, rag_files, attachments

# se vuoi forzare un base diverso
def _pick_generated_root() -> str:
    """
    Root di output per i file generati.
    Ordine preferenze:
    - env 'GENERATED_ROOT' (fallback al typo citato)
    - default: 'generated_<shortuuid>'
    """
    env = os.getenv("GENERATED_ROOT")
    if env:
        return env.rstrip("/")

    short = str(uuid.uuid4()).split("-")[0]
    return f"generated_{short}"

def _bucket_subdir(path: str) -> str:
    ext = (os.path.splitext(path)[1] or "").lower()
    if ext in CODE_EXTS:
        return "src"
    if ext in IMAGE_EXTS:
        return "images"
    if ext in DOC_EXTS:
        return "docs"
    return "docs"



def _retarget_files_under_generated(files: list[dict], prefix_path: str) -> list[dict]:
    """
    Riallincia i path dei file in base a GENERATED_ROOT e ai bucket {src, docs, images}.
    """
    log.info("_retarget_files_under_generated")
    base = _pick_generated_root()
    
    log.info("_retarget_files_under_generated --> Generated root:  %s", base)
    out: list[dict] = []
    temp_path =""
    for f in files or []:
        p = str(f.get("path") or "").lstrip("/").strip()
        
        c = f.get("content")
        if not p or c is None:
            continue
       
        sub = _bucket_subdir(p)
        
        intermediate_path = os.path.join(base, prefix_path)
        log.info(f"Retargeting 1 {p} to {intermediate_path} (subdir={sub})")
        # preserva solo il basename per evitare annidamenti sporchi
        bn = os.path.basename(p)
        log.info(f"Retargeting 2 bn to {bn} (subdir={sub})") 
        new_path = os.path.join(intermediate_path, sub, bn)
        out.append({"path": new_path, "content": c})
        log.info(f"Retargeting 3 {p} to {new_path}")
    return out


def _json_safe(obj):
    """Trasforma ricorsivamente set() -> list per garantire JSON serializzabile."""
    if isinstance(obj, dict):
        return {k: _json_safe(v) for k, v in obj.items()}
    if isinstance(obj, list):
        return [_json_safe(x) for x in obj]
    if isinstance(obj, set):
        return [_json_safe(x) for x in obj]  # list() di set
    return obj



def _inject_coding_system(msgs: list) -> list:
    """Garantisce un messaggio system che vieta prosa e impone il tool-call emit_files."""
    if msgs and isinstance(msgs[0], dict) and msgs[0].get("role") == "system" and "emit_files" in (msgs[0].get("content") or ""):
        return msgs
    sys = {
        "role": "system",
        "content": (
            "You are CLike code generator. You must produce output ONLY by CALLING the tool function "
            "'emit_files' with a JSON object: {\"files\":[{\"path\":\"<relative_path>\",\"content\":\"<file content>\"}]}. "
            "Never write normal assistant content. Do not include prose. If the user asks for code, return files via the tool call."
        ),
    }
    return [sys] + msgs

_CODE_FENCE_RE = re.compile(r"```(?P<lang>[a-zA-Z0-9+\-._]*)\s*\n(?P<code>.*?)(?:\r?\n)?```", re.DOTALL)

def _extract_files_from_fences(raw: str) -> list[dict]:
    files: list[dict] = []
    for i, m in enumerate(_CODE_FENCE_RE.finditer(raw or ""), start=1):
        lang = (m.group("lang") or "").strip().lower()
        code = m.group("code") or ""
        fname = _default_filename(lang, i)
        files.append({"path": fname, "content": code, "language": lang})
    return files

def _normalize_files_for_write(files: list[dict]) -> list[dict]:
    """
    Assicura path/content; se manca 'content' ma c'è 'text', usa text.
    Pulisce gli slash.
    """
    out: list[dict] = []
    for f in files or []:
        d = dict(f) if isinstance(f, dict) else {}
        path = (d.get("path") or "").strip()
        content = d.get("content")
        text = d.get("text")

        if (content is None or content == "") and isinstance(text, str) and text.strip():
            content = text
        if content is None:
            content = ""

        path = path.replace("\\", "/")
        out.append({"path": path, "content": content, "language": d.get("language","")})
    return out

def _extract_json(s: str) -> Dict[str, Any]:
    # 1) blocco ```json ... ```
    m = re.search(r"```json\s*(\{[\s\S]*?\})\s*```", s or "", re.M)
    if m:
        return json.loads(m.group(1))
    # 2) qualsiasi blocco ``` ... ``` con un oggetto json
    m = re.search(r"```\s*(\{[\s\S]*?\})\s*```", s or "", re.M)
    if m:
        return json.loads(m.group(1))
    # 3) fallback: primo { ... } nel testo
    i = (s or "").find("{"); j = (s or "").rfind("}")
    if i != -1 and j != -1 and j > i:
        return json.loads(s[i:j+1])
    raise ValueError("no valid JSON found")

def _default_filename(lang: str, idx: int = 1) -> str:
    l = (lang or "").lower()
    if l in ("py","python"): return f"module_{idx}.py"
    if l in ("ts","typescript"): return f"module_{idx}.ts"
    if l in ("js","javascript"): return f"module_{idx}.js"
    if l == "go": return f"module_{idx}.go"
    if l == "java": return f"module_{idx}.java"
    return f"module_{idx}.txt"

def _short_id(n: int = 8) -> str:
    return uuid.uuid4().hex[:n]

def _build_generation_roots(generation_id: str) -> Tuple[str, str, str, str]:
    """
    Ritorna (code_root_abs, test_root_abs, code_root_rel, test_root_rel)
    Esempio: ( .../src/generated_ab12cd34, .../tests/generated_ab12cd34, 'src/generated_ab12cd34', 'tests/generated_ab12cd34')
    """
    code_root_rel = os.path.join(settings.CODE_ROOT_BASE, f"{settings.GEN_ID_PREFIX}_{generation_id}")
    test_root_rel = os.path.join(settings.TEST_ROOT_BASE, f"{settings.GEN_ID_PREFIX}_{generation_id}")
    code_root_abs = os.path.join(settings.WORKSPACE_ROOT, code_root_rel)
    test_root_abs = os.path.join(settings.WORKSPACE_ROOT, test_root_rel)
    os.makedirs(os.path.join(code_root_abs, "src"), exist_ok=True)
    os.makedirs(os.path.join(code_root_abs, "doc"), exist_ok=True)
    os.makedirs(os.path.join(code_root_abs, "images"), exist_ok=True)
    os.makedirs(test_root_abs, exist_ok=True)
    return code_root_abs, test_root_abs, code_root_rel, test_root_rel

def _write_file_any(path: str, fobj: dict) -> None:
    os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
    if "content_base64" in fobj:
        data = base64.b64decode(fobj["content_base64"])
        with open(path, "wb") as wf:
            wf.write(data)
    else:
        content = fobj.get("content", "")
        with open(path, "w", encoding="utf-8") as wf:
            wf.write(content)

# ===== RAG hooks (best-effort; non bloccanti) =====
def _rag_project_id(body: dict) -> str:
    pid = (body or {}).get("project_id")
    if isinstance(pid, str) and pid.strip():
        return pid.strip()
    return "default"

RAG_TOP_K = int(os.getenv("RAG_TOP_K", "12"))

async def rag_index_items(project_id: str, items: list[dict]):
    # Optional server-side index; we prefer client-side, but keep for completeness.
    if not items:
        return
    payload = {"project_id": project_id, "items": []}
    for it in (items or []):
        p = (it.get("path") or "").strip()
        t = (it.get("text") or "").strip()
        if p and t:
            payload["items"].append({"path": p, "text": t})
    if not payload["items"]:
        return
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            await client.post(f"{_rag_base_url()}/index", json=payload)
    except Exception as e:
        log.warning("rag_index_items failed: %s", e)

async def rag_query(project_id: str, query: str, top_k: int = None):
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(f"{_rag_base_url()}/search",
                                  json={"project_id": project_id,
                                        "query": query or "",
                                        "top_k": int(top_k or RAG_TOP_K)})
            r.raise_for_status()
            data = r.json() or {}
            return data.get("hits") or []
    except Exception as e:
        log.warning("rag_query failed: %s", e)
        return []
# ----------------------------- models listing -------------------------------

def _mode_from_name(mid: str) -> str:
    if not mid:
        return "chat"
    low = mid.lower()
    if "embed" in low or "embedding" in low or "nomic-embed" in low:
        return "embed"
    return "chat"

def _normalize_models(payload: Dict[str, Any]) -> List[Dict[str, Any]]:
    # Case 1: CLike-style
    if isinstance(payload.get("models"), list):
        out = []
        for m in payload["models"]:
            name = m.get("name") or m.get("id") or m.get("model")
            if not name:
                continue
            mm = dict(m)
            mm.setdefault("name", name)
            mm.setdefault("modality", _mode_from_name(name))
            mm.setdefault("enabled", True)
            out.append(mm)
        return out
    # Case 2: OpenAI-style
    data = payload.get("data")
    if isinstance(data, list):
        out = []
        for m in data:
            mid = m.get("id")
            if not mid:
                continue
            out.append({
                "name": mid,
                "provider": "unknown",
                "modality": _mode_from_name(mid),
                "enabled": True,
                "capability": "medium",
                "latency": "medium",
                "cost": "medium",
                "privacy": "medium",
            })
        return out
    return []

def _filter_by_modality(models: List[Dict[str, Any]], modality: Optional[str]) -> List[Dict[str, Any]]:
    if modality == "chat":
        return [
            m for m in models
            if (m.get("modality") or "chat") in {"chat", "responses"}
        ]
    if modality in {"embed", "embeddings"}:
        return [
            m for m in models
            if (m.get("modality") or "chat") in {"embed", "embedding", "embeddings"}
        ]
    return models

async def _load_models_or_fallback() -> List[Dict[str, Any]]:
    # gateway
    try:
        base = str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/")
        async with httpx.AsyncClient(timeout=float(getattr(settings, "REQUEST_TIMEOUT_S", 60))) as client:
            r = await client.get(f"{base}/v1/models")
            r.raise_for_status()
            models = _normalize_models(r.json())
            if models:
                return models
    except Exception:
        pass
    # fallback: shared local catalog through llm_contracts
    try:
        catalog = await load_catalog(str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/"))
        raw = catalog.get("models", []) or []
        out = []
        for m in raw:
            name = m.get("name") or m.get("id") or m.get("model")
            if not name:
                continue
            mm = dict(m)
            mm.setdefault("name", name)
            mm.setdefault("modality", _mode_from_name(name))
            mm.setdefault("enabled", True)
            out.append(mm)
        return out
    except Exception:
        return []

async def _load_providers() -> Dict[str, Any]:
    """Provider availability snapshot from the gateway (single source of truth).

    On any failure we fall back to a permissive snapshot so a transient gateway
    hiccup never silently hides every cloud model.
    """
    permissive = {"providers": {}, "reasons": {}, "any_cloud": True, "any_local": True, "any": True}
    try:
        base = str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/")
        async with httpx.AsyncClient(timeout=float(getattr(settings, "REQUEST_TIMEOUT_S", 60))) as client:
            r = await client.get(f"{base}/v1/providers")
            r.raise_for_status()
            data = r.json()
            return data if isinstance(data, dict) else permissive
    except Exception:
        return permissive


_PROVIDER_KEY_ENV = {
    "openai": "OPENAI_API_KEY",
    "anthropic": "ANTHROPIC_API_KEY",
    "deepseek": "DEEPSEEK_API_KEY",
}


def _provider_available(provider: str, providers_snapshot: Dict[str, Any]) -> bool:
    prov = (provider or "").strip().lower()
    pmap = (providers_snapshot or {}).get("providers") or {}
    if prov in pmap:
        return bool(pmap[prov])
    # Unknown provider: do not block (matches gateway model_availability).
    return True


def _provider_unavailable_message(provider: str, model: str, providers_snapshot: Dict[str, Any]) -> str:
    prov = (provider or "").strip().lower()
    env = _PROVIDER_KEY_ENV.get(prov)
    if env:
        return (
            f"The selected model '{model}' uses the '{prov}' cloud provider, but its API key "
            f"({env}) is not configured. Set {env} to use cloud models, choose a model from a "
            f"configured provider, or switch Execution to 'agent only' to run locally."
        )
    reason = (providers_snapshot.get("reasons") or {}).get(prov) or f"the '{prov}' runtime is not reachable"
    return (
        f"The selected model '{model}' uses the local '{prov}' provider, but {reason}. Start the local "
        f"runtime, choose a model from a configured provider, or switch Execution to 'agent only' to run locally."
    )


def _provider_unavailable_envelope(provider: str, model: str, providers_snapshot: Dict[str, Any]) -> Dict[str, Any]:
    msg = _provider_unavailable_message(provider, model, providers_snapshot)
    return {
        "ok": False,
        "error_kind": "provider_unavailable",
        "provider": provider,
        "model": model,
        "text": msg,
        "errors": [msg],
        "files": [],
    }


# --- Local agent execution (free/coding) -------------------------------------
# Same architecture as Harper: the orchestrator owns prompt/context assembly and
# returns a package; the VS Code extension is the only place that spawns the
# codex/claude CLI. No cloud API key is required for this path.

_LOCAL_EXEC_PREFS = {"prefer_local_agent", "local_agent_only"}
_ROLE_LABEL = {"system": "System", "user": "User", "assistant": "Assistant"}

_CODING_LOCAL_SYSTEM = (
    "You are CLike, an expert engineer running locally inside the user's workspace. "
    "Generate exactly what the user asks (documentation, code, functions, images, etc.) "
    "as real files written under the directory '{output_root}/' (create it and any "
    "subdirectories as needed). Do not modify files outside '{output_root}/'. "
    "When finished, print a short plain-text summary of what you created."
)


def _execution_requests_local(value: Any) -> bool:
    return normalize_execution_preference(value) in _LOCAL_EXEC_PREFS


def _flatten_msgs_to_prompt(msgs: List[Dict[str, Any]]) -> str:
    parts = []
    for m in msgs or []:
        role = _ROLE_LABEL.get(str(m.get("role") or "user").strip().lower(), "User")
        content = str(m.get("content") or "").strip()
        if content:
            parts.append(f"{role}:\n{content}")
    return "\n\n".join(parts).strip()


def _local_execution_package(
    *, mode: str, prompt: str, executor_hint: Any, run_id: str, output_root: Optional[str] = None
) -> Dict[str, Any]:
    pkg: Dict[str, Any] = {
        "version": "1.0",
        "ok": True,
        "local_execution": True,
        "mode": mode,
        "prompt": prompt,
        "executor_hint": (str(executor_hint or "auto").strip().lower() or "auto"),
        "runId": run_id,
    }
    if output_root:
        pkg["output_root"] = output_root
    return pkg


@router.get("/models")
async def list_models(
    modality: Optional[str] = Query(default="chat", pattern="^(chat|embed|embeddings|all)$")
):
    try:
        models = await _load_models_or_fallback()
        providers = await _load_providers()

        # Show only enabled models in the UI/API list.
        models = [m for m in models if bool(m.get("enabled", True))]

        # Q1/Q2: hide models whose provider is not usable right now (missing key
        # for cloud, unreachable for local). Models without an explicit
        # `available` flag (e.g. degraded fallback catalog) are kept.
        models = [m for m in models if m.get("available", True)]

        if modality != "all":
            models = _filter_by_modality(models, modality)

        return {"version": "1.0", "models": models, "providers": providers}
    except Exception as ex:
        raise HTTPException(502, f"cannot load models: {type(ex).__name__}: {ex}")

# --------------------------------- Chat -------------------------------------

@router.post("/chat")
async def chat( req: Request):
    body = await req.json()
    mode = (body.get("mode") or "free" or "harper").lower()
    if mode not in ("free","harper"):
        raise HTTPException(400, "mode must be 'free' for /v1/chat")

    requested_provider = (body.get("provider") or "").lower().strip() or None
    requested_model = body.get("model") or "auto"

    llm_sel = await resolve_llm_selection(
        base_url=str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/"),
        mode="free",
        phase=None,
        requested_model=requested_model,
        requested_provider=requested_provider,
        profile_hint=body.get("profileHint"),
    )

    provider = llm_sel.get("provider") or ""
    model = llm_sel.get("model") or requested_model
    remote_name = llm_sel.get("remote_name") or model

    # Local-agent execution does not use the cloud at all, so the cloud provider
    # key check only applies when we are actually going to call the gateway (Q7).
    local_requested = _execution_requests_local(body.get("executionPreference"))
    if not local_requested:
        # Profile/routing resolution stays identical (Q7): if it lands on a provider
        # whose key is not configured, surface a clean message in the chat Text panel
        # instead of letting the gateway raise a raw 401.
        providers_snapshot = await _load_providers()
        if not _provider_available(provider, providers_snapshot):
            log.info("chat blocked: provider '%s' unavailable for model '%s'", provider, model)
            return _provider_unavailable_envelope(provider, model, providers_snapshot)

    mode_contract = llm_sel.get("mode_contract") or {}
    mode_contract = normalize_mode_contract(body.get("mode_contract"), mode_contract)
    try:
        validate_chat_contract(mode_contract)
    except Exception as e:
        raise HTTPException(400, f"invalid mode_contract for chat: {e}")
    messages = body.get("messages") or []
    if not isinstance(messages, list) or not messages:
        raise HTTPException(422, "messages (list) is required")

    # Attachments → inline vs rag
    attachments = body.get("attachments") or []
    # Normalize inputs: prefer explicit in_line_files/inline_files & rag_files.
    inline_files, rag_files, attachments = _normalize_context_from_body(body)
    log.info("chat extension & body fileds: %s, %s",  inline_files, rag_files)
    # If no explicit files were provided, but we have generic attachments, partition them.
    if not inline_files and not rag_files and attachments:
        inline_files, rag_files = await decide_inline_or_rag(attachments)
    
    log.info("chat attachments: %s, %s",  inline_files, rag_files)


    user_query = ""
    for m in reversed(messages):
        if (m.get("role") or "") == "user":
            user_query = (m.get("content") or "").strip()
            break

    # system + contesto
    sysmsg = {"role":"system","content":"You are CLike, a helpful and expert full-stack software engineering copilot."}
    msgs = [sysmsg] + list(messages)
    project_id = _rag_project_id(body)

    msgs = await _augment_messages_with_context(msgs, inline_files, rag_files, user_query, project_id)

    # RAG paths/inline opzionali (compat) TODO: the following code depends on evaluation if SPEC.md, IDEA.md or other file driven by VS extension are needed.
    # Legacy compatibility: only apply if NO new-style inline/rag files were provided
    if not inline_files and not rag_files:
        rag_paths  = body.get("rag_paths") or []
        rag_inline = body.get("rag_inline") or []
        if rag_paths or rag_inline:
            blobs = []
            if rag_paths:
                blobs.extend(_gather_rag_context(rag_paths))
            if rag_inline:
                blobs.extend([str(x) for x in rag_inline if x])
            if blobs:
                ctx = "\n\n".join(blobs[:8])
                msgs = [{"role":"system","content":"Use the following context if relevant:\n"+ctx}] + msgs

    # Local free (Q&A): hand the assembled prompt to the extension, which runs
    # the codex/claude CLI read-only and renders the answer in chat.
    if local_requested:
        prompt = _flatten_msgs_to_prompt(msgs)
        run_id = str(body.get("runId") or _short_id(8))
        log.info("chat local-execution package: mode=free run=%s executor=%s", run_id, body.get("localAgentExecutor"))
        return _local_execution_package(
            mode="free",
            prompt=prompt,
            executor_hint=body.get("localAgentExecutor"),
            run_id=run_id,
        )

    # validate modality
    all_models = await _load_models_or_fallback()
    requested_modality = next(
        (
            m.get("modality")
            for m in all_models
            if model in {
                m.get("id"),
                m.get("name"),
                m.get("remote_name"),
            }
        ),
        None,
    )
    if str(requested_modality or "").lower() in {"embed", "embedding", "embeddings"}:
        raise HTTPException(400, f"model '{model}' is an embedding model and cannot be used for chat.")
    
    log.info(
        "chat request: %s",
        json.dumps(
            {
                "requested_model": requested_model,
                "resolved_model": model,
                "remote_name": remote_name,
                "provider": provider,
                "profile": llm_sel.get("profile"),
                "messages_len": len(messages),
                "mode_contract": mode_contract,
            },
            ensure_ascii=False,
        ),
    )
    # Prepara meta per log
    _gw = str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/")
 
    # NOTE:
    # We intentionally do not build a local "payload" object here because /chat
    # sends arguments directly to call_gateway_chat(...).
    # Keep provider/model/profile as resolved by resolve_llm_selection().
    headers = {"Content-Type": "application/json"}
    _t0 = _time.time()
    try:
        all_models = await _load_models_or_fallback()
        model_entry = next((m for m in all_models if m.get("name") == model), None)
        req_max = int(body.get("max_tokens") or 2048)
        eff_max = su.tokens_per_model(msgs, model_entry, req_max)
        timeout_sec = min(340.0, 240.0 + (eff_max / 1000.0) * 3.8)
        
        text = await call_gateway_chat(
            model = model,
            messages = msgs,
            temperature= body.get("temperature"),
            max_tokens= eff_max,
            # --- AGGIUNGI: provider-awareness end-to-end ---        
            base_url= _gw, 
            timeout=timeout_sec,
            response_format=None, 
            tools=None, 
            tool_choice=None, 
            profile=llm_sel.get("profile"),
            provider=provider,
            mode_contract=mode_contract,

        )
        _ms = int((_time.time() - _t0) * 1000)
        log.info("chat response: %s", json.dumps({"text_len": len(text or ""), "latency_ms": _ms}, ensure_ascii=False))
        return {"version": "1.0", "text": text, "usage": {}, "sources": []}
    except Exception as e:
        _ms = int((_time.time() - _t0) * 1000)
        log.error("chat error: %s", json.dumps({"error": f"{type(e).__name__}: {e}", "latency_ms": _ms}, ensure_ascii=False))
        raise HTTPException(502, f"gateway chat failed: {type(e).__name__}: {e}")


def _gather_rag_context(paths: list[str], max_docs: int = 8, max_bytes: int = 200_000) -> list[str]:
    out = []
    for p in (paths or [])[:max_docs]:
        try:
            with open(p, "r", encoding="utf-8", errors="ignore") as f:
                t = f.read(max_bytes)
                if t and t.strip():
                    out.append(f"# Context: {p}\n{t.strip()}")
        except Exception:
            continue
    return out

def _kb(n_bytes: int) -> int:
    try: return int(n_bytes) // 1024
    except: return 0

def _fence(fname: str, content: str) -> str:
    lang = ""
    if fname.endswith(".py"): lang="python"
    elif fname.endswith(".ts"): lang="ts"
    elif fname.endswith(".js"): lang="js"
    elif fname.endswith(".go"): lang="go"
    elif fname.endswith(".java"): lang="java"
    return f"```{lang}\n# {fname}\n{content}\n```"
def _b64_to_bytes(b64: Optional[str]) -> Optional[bytes]:
    if not isinstance(b64, str) or not b64:
        return None
    try:
        # strip data URLs if present
        if b64.startswith("data:"):
            head, _, rest = b64.partition(",")
            b64 = rest
        return base64.b64decode(b64, validate=False)
    except Exception:
        return None

def _ext_from_path(p: str) -> str:
    try:
        return os.path.splitext((p or "").strip())[1].lower()
    except Exception:
        return ""
# --- ADD: extraction helpers ---
def _extract_text_from_pdf_bytes(raw: bytes) -> str:
    # pdfminer.six
    try:
        return extract_text(io.BytesIO(raw)) or ""
    except Exception as e:
        log.warning("PDF extract failed: %s", e)
        return ""
    


def _extract_text_from_docx_bytes(raw: bytes) -> str:
    # python-docx
    try:
        doc = docx.Document(io.BytesIO(raw))
        parts = []
        # paragraphs
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                parts.append(t)
        # tables (cells)
        for tbl in getattr(doc, "tables", []):
            for row in tbl.rows:
                for cell in row.cells:
                    t = (cell.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        log.warning("DOCX extract failed: %s", e)
        return ""
def _extract_text_from_xlsx_bytes(raw: bytes) -> str:
    if not openpyxl:
        log.warning("openpyxl non disponibile: skip xlsx")
        return ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
        parts = []
        for ws in wb.worksheets:
            parts.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                line = "\t".join(cells).strip()
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLSX extract failed: %s", e)
        return ""

def _extract_text_from_xls_bytes(raw: bytes) -> str:
    # Richiede xlrd>=2.0 (legge solo .xls)
    if not xlrd:
        log.warning("xlrd non disponibile: skip xls")
        return ""
    try:
        book = xlrd.open_workbook(file_contents=raw)
        parts = []
        for si in range(book.nsheets):
            sh = book.sheet_by_index(si)
            parts.append(f"# Sheet: {sh.name}")
            for r in range(sh.nrows):
                row = [str(sh.cell_value(r, c)) for c in range(sh.ncols)]
                line = "\t".join(row).strip()
                if line:
                    parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLS extract failed: %s", e)
        return ""

def _extract_text_from_xlsb_bytes(raw: bytes) -> str:
    if not open_xlsb:
        log.warning("pyxlsb non disponibile: skip xlsb")
        return ""
    try:
        parts = []
        with open_xlsb(io.BytesIO(raw)) as wb:
            for sheet_name in wb.sheets:
                parts.append(f"# Sheet: {sheet_name}")
                with wb.get_sheet(sheet_name) as sh:
                    for row in sh.rows():
                        vals = [str(c.v) if c.v is not None else "" for c in row]
                        line = "\t".join(vals).strip()
                        if line:
                            parts.append(line)
        return "\n".join(parts)
    except Exception as e:
        log.warning("XLSB extract failed: %s", e)
        return ""

def _extract_text_from_pptx_bytes(raw: bytes) -> str:
    if not Presentation:
        log.warning("python-pptx non disponibile: skip pptx")
        return ""
    try:
        prs = Presentation(io.BytesIO(raw))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"# Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text or "" for run in para.runs).strip()
                        if text:
                            parts.append(text)
                elif hasattr(shape, "text") and shape.text:
                    t = (shape.text or "").strip()
                    if t:
                        parts.append(t)
        return "\n".join(parts)
    except Exception as e:
        log.warning("PPTX extract failed: %s", e)
        return ""

async def decide_inline_or_rag(attachments: list[dict]) -> tuple[list[dict], list[dict]]:
    """
    Allineato all'estensione (partitionAttachments):
      - inline se a.content o a.bytes_b64 presenti
      - altrimenti, se a.path presente => RAG by path
      - altrimenti ignora (log warning)
    Niente soglie di size, niente budget qui (per coerenza end-to-end).
    """
    inline, rag = [], []
    if not attachments:
        return inline, rag

    for a in attachments:
        if not isinstance(a, dict):
            continue

        name   = a.get("name") or a.get("path") or "file"
        path   = a.get("path")
        origin = a.get("origin") or a.get("source")  # normalizza
        content    = a.get("content")
        bytes_b64  = a.get("bytes_b64")

        # Nota: evitiamo di loggare la base64 (solo boolean), per non intasare i log
        log.info("decide inline or rag: %s",
                 json.dumps({
                     "name": name,
                     "has_content": bool(content),
                     "has_bytes_b64": bool(bytes_b64),
                     "path": path,
                     "origin": origin
                 }, ensure_ascii=False))

        if content or bytes_b64:
            # Inline esattamente come fa l’estensione
            if bytes_b64:
                raw = _b64_to_bytes(bytes_b64)
                if raw:
                    ext = _ext_from_path(path)
                    if ext == ".pdf":
                        log.info("inline: PDF")
                        txt = _extract_text_from_pdf_bytes(raw)
                    elif ext == ".docx":
                        log.info("inline DOCX")
                        txt = _extract_text_from_docx_bytes(raw)
                    elif ext == ".xlsx":
                        log.info("inline: XLSX")
                        txt = _extract_text_from_xlsx_bytes(raw)
                    elif ext == ".xls":
                        log.info("inline: XLS (legacy)")
                        txt = _extract_text_from_xls_bytes(raw)
                    elif ext == ".xlsb":
                        log.info("inline: XLSB")
                        txt = _extract_text_from_xlsb_bytes(raw)
                    elif ext == ".pptx":
                        log.info("inline: PPTX")
                        txt = _extract_text_from_pptx_bytes(raw)
                    else:
                        # fallback: se è testo “grezzo” o sconosciuto, prova a decodare come utf-8
                        try:
                            txt = raw.decode("utf-8", errors="ignore")
                        except Exception:
                            txt = ""
                  
                    log.info("inline file %s -> %d chars", path, len(txt))
                    content = txt

            inline.append({
                "name": name,
                "path": path,          # opzionale (può servire per tracciabilità)
                "content": content,    # può essere None
                "bytes_b64": bytes_b64,# può essere None
                "origin": origin
            })
        elif path:
            # RAG by path, minimale (non inoltriamo bytes_b64 per non gonfiare la payload)
            rag.append({
                "name": name,
                "path": path,
                "origin": origin
            })
        else:
            log.warning("Attachment senza content/bytes_b64 e senza path: ignorato: %s", name)

    log.info("attachments routing %s",
             json.dumps({"inline": len(inline), "rag": len(rag)}, ensure_ascii=False))
    return inline, rag


async def _augment_messages_with_context(
    msgs: list[dict],
    inline_files: list[dict],
    rag_files: list[dict],
    user_query: str,
    project_id: str
) -> list[dict]:
    """
    Enrich messages with:
      1) Inline files as fenced blocks.
      2) Exact-path RAG fetch for explicitly attached rag_files.
      3) Semantic RAG search only when no exact rag paths are available.
    """
    out = list(msgs)

    # 1) Inline files
    if inline_files:
        blocks = "\n\n".join(
            _fence(f.get("name") or f.get("path") or "file", f.get("content") or "")
            for f in inline_files
            if isinstance(f, dict) and (f.get("content") or "").strip()
        )
        if blocks.strip():
            out = [{
                "role": "system",
                "content": "You can use the following inline project files:\n\n" + blocks
            }] + out

    rag_paths = []
    seen_paths = set()
    for rf in (rag_files or []):
        if not isinstance(rf, dict):
            continue
        p = (rf.get("path") or "").strip()
        if not p:
            continue
        norm = os.path.normpath(p).replace("\\", "/").lower()
        if norm in seen_paths:
            continue
        seen_paths.add(norm)
        rag_paths.append(p)

    # 2) Exact-path fetch for explicitly attached RAG files
    if rag_paths:
        try:
            store = RagStore(project_id=project_id)
            docs = await store.fetch_docs_by_paths(
                rag_paths,
                max_chars_per_doc=12000,
                limit_points=max(200, len(rag_paths) * 50),
            )

            if docs:
                blocks = []
                for d in docs:
                    text = (d.get("text") or "").strip()
                    path = (d.get("path") or "").strip()
                    if not text or not path:
                        continue
                    blocks.append(f"### {path}\n{text}")

                if blocks:
                    out = [{
                        "role": "system",
                        "content": "Attached project context:\n\n" + "\n\n".join(blocks)
                    }] + out
                    return out
        except Exception as e:
            log.warning("Exact-path RAG enrichment failed: %s", e)

        # If explicit rag paths were provided but exact fetch produced nothing,
        # do not silently switch to semantic search over arbitrary workspace content.
        return out

    # 3) Semantic RAG only when the user did not attach explicit rag paths
    try:
        q = (user_query or "").strip()
        if q:
            hits = await rag_query(project_id, q, top_k=RAG_TOP_K)
            dedup = []
            seen = set()

            for h in hits or []:
                if not isinstance(h, dict):
                    continue
                path = (h.get("path") or "").strip()
                text = (h.get("text") or "").strip()
                chunk = int(h.get("chunk", 0))
                if not path or not text:
                    continue

                sig = (path.lower(), chunk, text[:96])
                if sig in seen:
                    continue
                seen.add(sig)
                dedup.append({"path": path, "chunk": chunk, "text": text})

            if dedup:
                blocks = []
                for h in dedup[:RAG_TOP_K]:
                    t = h["text"]
                    if len(t) > 4000:
                        t = t[:4000] + "\n...[truncated]..."
                    blocks.append(f"### {h['path']}:{h['chunk']}\n{t}")

                if blocks:
                    out = [{
                        "role": "system",
                        "content": "Relevant project context:\n\n" + "\n\n".join(blocks)
                    }] + out
    except Exception as e:
        log.warning("Semantic RAG enrichment failed: %s", e)

    return out

@router.post("/generate")
async def generate(req: Request):
    log.info("generate request +++")
    body = await req.json()
    mode = (body.get("mode") or "coding").lower()
    if mode not in ("harper", "coding"):
        raise HTTPException(400, "mode must be 'harper' or 'coding'")

    requested_model = body.get("model") or "auto"
    requested_provider = (body.get("provider") or "").lower().strip() or None

    llm_sel = await resolve_llm_selection(
        base_url=str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/"),
        mode="coding" if mode == "coding" else "harper",
        phase="kit" if mode == "coding" else body.get("phase"),
        requested_model=requested_model,
        requested_provider=requested_provider,
        profile_hint=body.get("profileHint"),
    )

    model = llm_sel.get("model") or requested_model
    provider = (llm_sel.get("provider") or "").lower().strip()
    remote_name = llm_sel.get("remote_name") or model

    local_requested = _execution_requests_local(body.get("executionPreference"))
    if not local_requested:
        providers_snapshot = await _load_providers()
        if not _provider_available(provider, providers_snapshot):
            log.info("generate blocked: provider '%s' unavailable for model '%s'", provider, model)
            return _provider_unavailable_envelope(provider, model, providers_snapshot)

    mode_contract = llm_sel.get("mode_contract") or {}
    mode_contract = normalize_mode_contract(body.get("mode_contract"), mode_contract)

    PROVIDERS_RESPONSE_FORMAT = {"openai", "azure_openai"}
    PROVIDERS_TOOL_CALL = {"ollama", "anthropic", "deepseek", "vllm"}

    prov = (provider or "").lower().strip()
    _use_respfmt = prov in PROVIDERS_RESPONSE_FORMAT
    _use_tools = prov in PROVIDERS_TOOL_CALL

    messages = body.get("messages") or []
    # Enforce tool-call in coding mode
    messages = _inject_coding_system(messages)

    if not isinstance(messages, list) or not messages:
        raise HTTPException(422, "messages (list) is required")

    # generation id + roots (path pianificati; nessuna scrittura ancora)
    gen_id = _short_id(8)
    _code_abs, _test_abs, code_root_rel, _test_rel = _build_generation_roots(gen_id)

    # System che “inchioda” lo schema di uscita (usato come contesto, ma non forziamo più response_format qui)
    sys_schema = {
        "role": "system",
        "content": (
            "You are CLike an expert code generator, Image and Video creator, UI/UX desinger with Cloud Skills, Application and Infrastructure Architect and more.  ALWAYS answer ONLY valid JSON with this schema:\n"
            "{\n"
            '  "files": [ { "path": "<relative/path/with/extension>", "content": "<full file content>" } ],\n'
            '  "messages": [ { "role": "assistant", "content": "<optional explanation>" } ]\n'
            "}\n"
            "No code fences. No 'Generated files:' lists. If multiple languages are needed, include multiple entries in files[].\n"
        )
    }

    # Attachments → inline vs rag
    attachments = body.get("attachments") or []
    # Normalize inputs: prefer explicit in_line_files/inline_files & rag_files.
    inline_files, rag_files, attachments = _normalize_context_from_body(body)

    # If no explicit files were provided, but we have generic attachments, partition them.
    if not inline_files and not rag_files and attachments:
        inline_files, rag_files = await _decide_inline_or_rag(attachments)
    
    # RAG query dall’ultimo user
    user_query = ""
    for m in reversed(messages):
        if (m.get("role") or "") == "user":
            user_query = (m.get("content") or "").strip()
            break

    msgs = [sys_schema] + list(messages)
    project_id = _rag_project_id(body)
    msgs = await _augment_messages_with_context(msgs, inline_files, rag_files, user_query, project_id)

    # Local coding: instead of asking the model for JSON files[], instruct the
    # local agent to write real files under a generated root, then let the
    # extension collect them and render the synthesis in chat + Files tab.
    if local_requested:
        output_root = f"generated/{gen_id}"
        write_sys = {"role": "system", "content": _CODING_LOCAL_SYSTEM.format(output_root=output_root)}
        local_msgs = [write_sys] + list(messages)
        local_msgs = await _augment_messages_with_context(local_msgs, inline_files, rag_files, user_query, project_id)
        prompt = _flatten_msgs_to_prompt(local_msgs)
        run_id = str(body.get("runId") or gen_id)
        log.info("generate local-execution package: mode=coding run=%s root=%s executor=%s", run_id, output_root, body.get("localAgentExecutor"))
        return _local_execution_package(
            mode="coding",
            prompt=prompt,
            executor_hint=body.get("localAgentExecutor"),
            run_id=run_id,
            output_root=output_root,
        )

    # modality check
    all_models = await _load_models_or_fallback()
    requested_modality = next(
        (
            m.get("modality")
            for m in all_models
            if model in {
                m.get("id"),
                m.get("name"),
                m.get("remote_name"),
            }
        ),
        None,
    )
    if str(requested_modality or "").lower() in {"embed", "embedding", "embeddings"}:
        raise HTTPException(400, f"model '{model}' is an embedding model and cannot be used for code generation.")

    log.info("generate request: %s", json.dumps({"model": model, "messages_len": len(messages)}, ensure_ascii=False))

    # ======== Chiamata gateway (prima scelta: TOOL CALLING) ========
    # ======== Gateway payload builder (coding) ========
    base_url = str(getattr(settings, "GATEWAY_URL", "http://localhost:8000")).rstrip("/")

   
    FILES_BUNDLE_SCHEMA = {
        "name": "files_bundle_v1",
        "schema": {
            "type": "object",
            "properties": {
                "files": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string"},
                            "content": {"type": "string"},
                        },
                        "required": ["path", "content"],
                        "additionalProperties": False,
                    }
                }
            },
            "required": ["files"],
            "additionalProperties": False,
        },
        "strict": True,
    }

    emit_files_guidance = (
        "When you propose code changes, you MUST return files, not prose. "
        "If tool calling is available, you MUST call emit_files with a non-empty files array. "
        "The tool arguments must be exactly: {\"files\":[{\"path\":\"...\",\"content\":\"...\"}]}. "
        "Do not call emit_files with empty input. "
        "If tool calling is not available, return a single top-level JSON object with the structure "
        "{\"files\":[{\"path\":\"...\",\"content\":\"...\"}]} and no extra text before or after. "
        "Do not return explanations outside the file bundle."
    )
    
    msgs = [{"role": "system", "content": emit_files_guidance}] + msgs

    temperature = body.get("temperature", 0.1)
    max_tokens = body.get("max_tokens", 4048)

    payload = {
        "model": model,
        "messages": msgs,
        "base_url": base_url,
        "remote_name": remote_name,
        "profile": llm_sel.get("profile"),
        "mode_contract": mode_contract,
    }

    if provider:
        payload["provider"] = provider

    if str(model).startswith("gpt-5"):
        if max_tokens is not None:
            payload["max_completion_tokens"] = max_tokens
    else:
        if max_tokens is not None:
            payload["max_tokens"] = max_tokens

    if temperature is not None and not str(model).startswith("gpt-5"):
        payload["temperature"] = temperature

    emit_files_tool = {
        "type": "function",
        "function": {
            "name": "emit_files",
            "description": "Return source files to be written by the caller.",
            "parameters": {
                "type": "object",
                "properties": {
                    "files": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "path": {"type": "string"},
                                "content": {"type": "string"},
                                "language": {"type": "string"},
                                "executable": {"type": "boolean"},
                            },
                            "required": ["path", "content"],
                            "additionalProperties": False,
                        },
                    }
                },
                "required": ["files"],
                "additionalProperties": False,
            },
        },
    }

    try:
        payload = apply_generate_contract(
            payload=payload,
            provider=provider,
            contract=mode_contract,
            files_bundle_schema=FILES_BUNDLE_SCHEMA,
            emit_files_tool=emit_files_tool,
        )
    except Exception as e:
        raise HTTPException(400, f"invalid mode_contract for generate: {e}")

    payload = _json_safe(payload)

    _headers = {"Content-Type": "application/json", "X-CLike-Profile": llm_sel.get("profile") or "code.strict"}
    if provider:
        _headers["X-CLike-Provider"] = provider

    log.info(
        "gateway.request %s",
        json.dumps(
            {
                "url": f"{base_url}/v1/chat/completions",
                "model": model,
                "profile": payload.get("profile"),
                "tools": bool(payload.get("tools")),
                "tool_choice": bool(payload.get("tool_choice")),
                "has_response_format": bool(payload.get("response_format")),
                "max_tokens": payload.get("max_tokens"),
                "max_completion_tokens": payload.get("max_completion_tokens"),
                "provider": payload.get("provider"),
                "mode_contract": mode_contract,
            },
            ensure_ascii=False,
        ),
    )
    try:
        all_models = await _load_models_or_fallback()
        model_entry = next((m for m in all_models if m.get("name") == model), None)
        req_max = int(body.get("max_tokens") or 2048)
        eff_max = su.tokens_per_model(msgs, model_entry, req_max)
        timeout_sec = min(440.0, 400.0 + (eff_max / 1000.0) * 4.0)
        payload["timeout"] = timeout_sec

        data = await call_gateway_generate(payload, _headers)

        def _is_truncated_no_files_response(d: Any) -> bool:
            if not isinstance(d, dict):
                return False

            finish_reason = str(d.get("finish_reason") or "").lower()
            raw = d.get("raw") or {}

            raw_stop_reason = str(raw.get("stop_reason") or "").lower()
            incomplete_reason = str((raw.get("incomplete_details") or {}).get("reason") or "").lower()

            if finish_reason in {"length", "max_tokens", "max_output_tokens"}:
                return True
            if raw_stop_reason in {"max_tokens"}:
                return True
            if incomplete_reason in {"max_output_tokens", "max_tokens"}:
                return True

            return False
        
        # Single retry for truncated coding generations.
        # This helps both:
        # - GPT-5 family producing incomplete textual JSON
        # - Claude stopping during emit_files input serialization
        if _is_truncated_no_files_response(data):
            retry_payload = dict(payload)

            if "max_completion_tokens" in retry_payload:
                retry_payload["max_completion_tokens"] = max(int(retry_payload["max_completion_tokens"]) * 2, 12000)

            if "max_tokens" in retry_payload:
                retry_payload["max_tokens"] = max(int(retry_payload["max_tokens"]) * 2, 12000)

            retry_payload["timeout"] = max(float(payload.get("timeout") or 0), 520.0)

            log.info(
                "generate retry-on-truncation %s",
                json.dumps(
                    {
                        "provider": provider,
                        "model": model,
                        "retry_max_tokens": retry_payload.get("max_tokens"),
                        "retry_max_completion_tokens": retry_payload.get("max_completion_tokens"),
                    },
                    ensure_ascii=False,
                ),
            )

            data_retry = await call_gateway_generate(retry_payload, _headers)

            # Prefer retry result if it is still truncated but at least different / fuller.
            if isinstance(data_retry, dict):
                data = data_retry
        def _has_empty_emit_files_tool_call(d: Any) -> bool:
            if not isinstance(d, dict):
                return False

            raw = d.get("raw") or {}
            content = raw.get("content") or []
            if not isinstance(content, list):
                return False

            for item in content:
                if not isinstance(item, dict):
                    continue
                if item.get("type") == "tool_use" and item.get("name") == "emit_files":
                    tool_input = item.get("input")
                    if tool_input == {} or tool_input is None:
                        return True
            return False

        # ---- Estrazione FILES in modo robusto cross-provider ----
        files: List[Dict[str, Any]] = []

        def _first_message(d: Any) -> Dict[str, Any]:
            """Ritorna in sicurezza il primo message da data['choices'][0]['message'] se presente e ben formato."""
            try:
                if isinstance(d, dict):
                    ch = d.get("choices")
                    if isinstance(ch, list) and ch:
                        c0 = ch[0]
                        if isinstance(c0, dict):
                            m = c0.get("message")
                            if isinstance(m, dict):
                                return m
            except Exception:
                pass
            return {}

        msg = _first_message(data)               # dict (o {})
        content_str = ""
        if isinstance(msg, dict):
            content_str = msg.get("content") or ""

        # 1) Preferisci tool_calls (OpenAI compat)
        tool_calls = msg.get("tool_calls") if isinstance(msg, dict) else None
        if isinstance(tool_calls, list) and tool_calls:
            for tc in tool_calls:
                try:
                    if (tc.get("type") == "function") and (tc.get("function", {}).get("name") == "emit_files"):
                        args_raw = tc.get("function", {}).get("arguments")
                        args = json.loads(args_raw) if isinstance(args_raw, str) else (args_raw or {})
                        parsed = (args.get("files") or [])
                        files = _normalize_files_for_write(parsed)
                        log.info("generate tool_calls->files %s", json.dumps({"count": len(files)}, ensure_ascii=False))
                        break
                except Exception:
                    # continua coi fallback
                    pass

        # 2) Fallback: top-level "files" (es. Ollama / adapter custom)
        if not files and isinstance(data, dict) and isinstance(data.get("files"), list):
            files = _normalize_files_for_write(data["files"])
            log.info("generate top-level->files %s", json.dumps({"count": len(files)}, ensure_ascii=False))
        if not files and isinstance(data, dict):
            maybe_text = data.get("text")
            if isinstance(maybe_text, str) and maybe_text.strip():
                try:
                    obj = _extract_json(maybe_text)
                    jf = obj.get("files") if isinstance(obj, dict) else None
                    if isinstance(jf, list) and jf:
                        files = _normalize_files_for_write(jf)
                        log.info("generate top-level-text-json->files %s", json.dumps({"count": len(files)}, ensure_ascii=False))
                except Exception:
                    pass
        # 3) Fallback: JSON puro dentro message.content / oppure 'text'
        if not files:
            if not content_str and isinstance(data, dict) and "text" in data:
                content_str = data.get("text") or ""

            if isinstance(content_str, str) and content_str.strip():
                try:
                    obj = _extract_json(content_str)
                    jf = obj.get("files") if isinstance(obj, dict) else None
                    if isinstance(jf, list) and jf:
                        files = _normalize_files_for_write(jf)
                        log.info("generate content-json->files %s", json.dumps({"count": len(files)}, ensure_ascii=False))
                except Exception:
                    pass

                # Final tolerant fallback: any fenced code block becomes a file.
                if not files:
                    from_fences = _extract_files_from_fences(content_str)
                    if from_fences:
                        files = _normalize_files_for_write(from_fences)
                        log.info("generate fences->files %s", json.dumps({"count": len(files)}, ensure_ascii=False))

        log.info("generate files (post-extract) %s", json.dumps({"count": len(files)}, ensure_ascii=False))

        # 5) Se ancora vuoto → 422 coerente
        if not files:
            log.info("generate no-files (nothing from tool_calls/top-level/json/fences)")
            raise HTTPException(status_code=422, detail="model did not produce 'files' with path+content")

        # 6) retarget sotto generated_<uuid>/ {src|docs|images}
        temp_path = str(uuid.uuid4()).split("-")[0]
        files = _retarget_files_under_generated(files, temp_path)   # <— prima dei diff!

        # 7) diffs
        diffs: List[Dict[str, Any]] = []
        for fobj in files:
            path = fobj["path"]
            content = fobj.get("content", "")
            prev = su.read_file(path) or ""
            patch = su.to_diff(prev, content, path)
            diffs.append({"path": path, "diff": patch})

        # 8) risposta completa (popola "text" e "diffs" per i tab)
        result = {
            "version": "1.0",
            "files": files,
            "usage": (data.get("usage") if isinstance(data, dict) else {}) or {},
            "sources": [],
            "text": "Generated files:\n" + "\n".join(f"- {f['path']}" for f in files),
            "diffs": diffs or ["(No diffs computed: new files)"],
            "audit_id": "coding-toolcalls",
        }
        return result

    except httpx.HTTPStatusError as e:
        # Propaga il vero body (niente 502 generici)
        raise HTTPException(e.response.status_code, detail=f"gateway chat failed: {e.response.text}")
    except Exception as e:
        raise HTTPException(502, f"gateway chat failed: {type(e).__name__}: {e}")

    
   

# -------------------------------- Apply -------------------------------------

@router.post("/apply")
async def apply(req: Request):
    """
    Applica file **direttamente dal payload**:
      {
        "files": [{ "path":"...", "content":"..." }, ...],
        "selection": { "apply_all": true }    # oppure: { "paths": ["a","b"] }
      }

    Nota: supporto a run_dir è stato rimosso.
    """
    body = await req.json()

    # rifiuta legacy
    if body.get("run_dir"):
        raise HTTPException(400, "run_dir is no longer supported. Pass 'files' directly in the request body.")

    files = body.get("files")
    if not isinstance(files, list) or not files:
        raise HTTPException(400, "files (list) is required")

    selection = body.get("selection") or {}
    paths_selected: set[str] = set()
    if isinstance(selection, dict):
        if selection.get("apply_all"):
            paths_selected = { (f.get("path") or "").strip() for f in files if isinstance(f, dict) }
        else:
            for p in selection.get("paths", []):
                if isinstance(p, str) and p.strip():
                    paths_selected.add(p.strip())

    applied: list[str] = []
    failures: list[dict] = []

    for fobj in files:
        if not isinstance(fobj, dict):
            continue
        path = (fobj.get("path") or "").strip()
        if not path:
            continue
        if paths_selected and path not in paths_selected:
            continue
        try:
            _write_file_any(path, fobj)
            applied.append(path)
        except Exception as e:
            failures.append({"path": path, "error": f"{type(e).__name__}: {e}"})

    log.info("apply result: %s", json.dumps({"applied": len(applied), "failures": len(failures)}, ensure_ascii=False))

    if failures:
        return {"applied": applied, "failures": failures}
    return {"applied": applied}
